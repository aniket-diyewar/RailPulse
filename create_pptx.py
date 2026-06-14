from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

DARK_BG = RGBColor(0x0A, 0x0F, 0x1A)
CARD_BG = RGBColor(0x1E, 0x29, 0x3B)
BLUE_GRAD = RGBColor(0x60, 0xA5, 0xFA)
PURPLE_GRAD = RGBColor(0xA7, 0x8B, 0xFA)
GREEN = RGBColor(0x4A, 0xDE, 0x80)
YELLOW = RGBColor(0xFB, 0xBF, 0x24)
RED = RGBColor(0xEF, 0x44, 0x44)
TEXT_PRIMARY = RGBColor(0xF1, 0xF5, 0xF9)
TEXT_SECONDARY = RGBColor(0x94, 0xA3, 0xB8)
TEXT_MUTED = RGBColor(0x64, 0x74, 0x8B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, color, opacity=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    tf = shape.text_frame
    tf.word_wrap = True
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False, color=TEXT_PRIMARY, align=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_rich_text_box(slide, left, top, width, height, lines, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, size, bold, color) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = Pt(6)
    return txBox

def add_bullet_box(slide, left, top, width, height, items, font_size=18, color=TEXT_SECONDARY, bullet="•"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"{bullet} {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(6)
    return txBox

def add_line(slide, left, top, width, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

# ── Slide 1: Title ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_shape(slide, Inches(0.8), Inches(1.8), Inches(11.7), Inches(4.2), CARD_BG)
add_line(slide, Inches(0.8), Inches(1.8), Inches(11.7), BLUE_GRAD)
add_text_box(slide, Inches(2), Inches(2.2), Inches(9.3), Inches(1.5),
    "RailPulse", font_size=60, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(slide, Inches(2), Inches(3.5), Inches(9.3), Inches(1.2),
    "Railway Intelligence System", font_size=32, color=BLUE_GRAD, align=PP_ALIGN.CENTER)
add_line(slide, Inches(5), Inches(4.6), Inches(3.3), PURPLE_GRAD)
add_text_box(slide, Inches(2), Inches(4.8), Inches(9.3), Inches(0.8),
    "Real-Time Monitoring | ML Delay Prediction | Analytics Dashboard", font_size=18, color=TEXT_SECONDARY, align=PP_ALIGN.CENTER)
add_text_box(slide, Inches(2), Inches(5.5), Inches(9.3), Inches(0.5),
    "Built for Railways Hackathon", font_size=16, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

# ── Slide 2: Problem ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_text_box(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.8),
    "The Problem", font_size=40, bold=True, color=WHITE)
add_line(slide, Inches(0.8), Inches(1.2), Inches(3), BLUE_GRAD)
add_bullet_box(slide, Inches(1), Inches(1.8), Inches(11), Inches(5), [
    "Railway operators rely on manual monitoring — no unified live view",
    "Delay prediction is reactive, not proactive",
    "Historical data scattered across silos, no central analytics",
    "No real-time public dashboard for train performance",
    "Operators lack tools to identify systemic delay patterns",
], font_size=20, color=TEXT_SECONDARY)

# ── Slide 3: Solution ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_text_box(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.8),
    "Our Solution", font_size=40, bold=True, color=WHITE)
add_line(slide, Inches(0.8), Inches(1.2), Inches(3), BLUE_GRAD)
items = [
    ("AI-Powered Dashboard", "Unified live view of train statuses with real-time NTES feed"),
    ("ML Delay Prediction", "Random Forest model predicts delays with 74% accuracy"),
    ("Interactive Analytics", "Station, weather, peak-hour, route, and monthly trend analysis"),
    ("Live + Fallback", "Primary NTES API data with automatic simulation fallback"),
    ("Nagpur-Focused", "Zeroed in on Nagpur Junction for faster, reliable live data"),
]
for i, (title, desc) in enumerate(items):
    y = Inches(1.6) + Inches(1.1) * i
    card = add_shape(slide, Inches(1), y, Inches(11.3), Inches(0.95), CARD_BG)
    add_text_box(slide, Inches(1.3), y + Inches(0.08), Inches(10.7), Inches(0.4),
        title, font_size=22, bold=True, color=BLUE_GRAD)
    add_text_box(slide, Inches(1.3), y + Inches(0.48), Inches(10.7), Inches(0.4),
        desc, font_size=16, color=TEXT_SECONDARY)

# ── Slide 4: Tech Stack ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_text_box(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.8),
    "Technology Stack", font_size=40, bold=True, color=WHITE)
add_line(slide, Inches(0.8), Inches(1.2), Inches(3), BLUE_GRAD)

stack = [
    ("Frontend", "Streamlit", "Python-based reactive UI framework"),
    ("Machine Learning", "Scikit-learn", "Random Forest (74% acc, 7.5 MAE)"),
    ("Data Processing", "Pandas, NumPy, PyArrow", "Parquet storage, 100k+ records"),
    ("Visualization", "Plotly, Folium", "Interactive charts + route network map"),
    ("Real-Time Data", "NTES API (ntes-client)", "Live Indian Railways station feed"),
    ("Data Generation", "Synthetic Generator", "30 stations, 46 routes, 3-year span"),
]
for i, (layer, tech, desc) in enumerate(stack):
    col = i % 3
    row = i // 3
    x = Inches(0.8) + Inches(4.1) * col
    y = Inches(1.6) + Inches(2.8) * row
    card = add_shape(slide, x, y, Inches(3.8), Inches(2.5), CARD_BG)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.15), Inches(3.4), Inches(0.4),
        layer, font_size=18, bold=True, color=PURPLE_GRAD)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.6), Inches(3.4), Inches(0.3),
        tech, font_size=22, bold=True, color=WHITE)
    add_text_box(slide, x + Inches(0.2), y + Inches(1.1), Inches(3.4), Inches(1.2),
        desc, font_size=14, color=TEXT_SECONDARY)

# ── Slide 5: Live Dashboard ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
    "Live Dashboard — Nagpur Junction", font_size=40, bold=True, color=WHITE)
add_line(slide, Inches(0.8), Inches(1.2), Inches(3), BLUE_GRAD)
add_text_box(slide, Inches(1), Inches(1.8), Inches(11.3), Inches(0.5),
    "Real-time train departures from Nagpur (NGP) via NTES API", font_size=20, color=TEXT_SECONDARY)
add_shape(slide, Inches(1), Inches(2.5), Inches(11.3), Inches(0.7), CARD_BG)
add_text_box(slide, Inches(1.3), Inches(2.55), Inches(10.7), Inches(0.6),
    "4 Active Trains  •  100% On-Time  •  0.0 min Avg Delay  •  0% Delayed", font_size=22, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

trains = [
    ("NGP VANDE BHARAT", "00123", "NAGPUR → MUMBAI", "19:00", "ON TIME"),
    ("SEWAGRAM EXP", "12140", "NAGPUR → CSMT", "21:15", "ON TIME"),
    ("GITANJALI EXP", "12859", "CSMT → HOWRAH", "19:00", "ON TIME"),
    ("KERALA SF EXP", "12625", "TVC → NEW DELHI", "12:10", "DELAYED 5 min"),
]
for i, (name, t_id, route, time, status) in enumerate(trains):
    y = Inches(3.5) + Inches(0.9) * i
    color = GREEN if "ON TIME" in status else YELLOW
    card = add_shape(slide, Inches(1), y, Inches(11.3), Inches(0.75), CARD_BG)
    add_text_box(slide, Inches(1.3), y + Inches(0.05), Inches(5), Inches(0.3),
        f"{name}  #{t_id}", font_size=18, bold=True, color=WHITE)
    add_text_box(slide, Inches(1.3), y + Inches(0.38), Inches(5), Inches(0.3),
        route, font_size=14, color=TEXT_SECONDARY)
    add_text_box(slide, Inches(7.5), y + Inches(0.05), Inches(2), Inches(0.3),
        f"🕐 {time}", font_size=16, color=TEXT_MUTED)
    add_text_box(slide, Inches(9.5), y + Inches(0.15), Inches(2.5), Inches(0.4),
        status, font_size=16, bold=True, color=color, align=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(6.8), Inches(11.3), Inches(0.4),
    "🔴 LIVE  |  Updated 19:42:15  |  Data via NTES Indian Railways", font_size=14, color=TEXT_MUTED, align=PP_ALIGN.RIGHT)

# ── Slide 6: ML Model ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
    "ML Delay Prediction Model", font_size=40, bold=True, color=WHITE)
add_line(slide, Inches(0.8), Inches(1.2), Inches(3), BLUE_GRAD)

stats = [
    ("74.6%", "Classification Accuracy", GREEN),
    ("7.5 min", "Mean Absolute Error", YELLOW),
    ("13", "Features", BLUE_GRAD),
    ("100k+", "Training Records", PURPLE_GRAD),
]
for i, (val, label, color) in enumerate(stats):
    x = Inches(0.8) + Inches(3.1) * i
    card = add_shape(slide, x, Inches(1.6), Inches(2.8), Inches(2.2), CARD_BG)
    add_text_box(slide, x, Inches(1.8), Inches(2.8), Inches(1), val, font_size=48, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text_box(slide, x, Inches(2.9), Inches(2.8), Inches(0.5), label, font_size=16, color=TEXT_SECONDARY, align=PP_ALIGN.CENTER)

add_text_box(slide, Inches(0.8), Inches(4.2), Inches(11.7), Inches(0.5),
    "Model Architecture", font_size=24, bold=True, color=WHITE)
add_bullet_box(slide, Inches(1), Inches(4.7), Inches(11.3), Inches(2.5), [
    "Algorithm: Random Forest Regressor (delay in min) + Random Forest Classifier (on-time vs delayed)",
    "Features: hour, day_of_week, month, is_peak_hour, weather, congestion_index, distance_km, train_name, from_station, to_station",
    "Train/Test split: 80/20 with stratified sampling for classifier",
    "Models compressed with joblib (compress=3) — 47 MB → 13 MB",
    "Predictor fetches trained encoders + scalers for real-time inference",
], font_size=17, color=TEXT_SECONDARY)

# ── Slide 7: Data Flow ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
    "Data Flow Architecture", font_size=40, bold=True, color=WHITE)
add_line(slide, Inches(0.8), Inches(1.2), Inches(3), BLUE_GRAD)

boxes = [
    (Inches(0.8), Inches(1.8), "NTES API\nIndian Railways", BLUE_GRAD),
    (Inches(3.8), Inches(1.8), "RailwayFetcher\nNGP Station", PURPLE_GRAD),
    (Inches(6.8), Inches(1.8), "Cache Layer\n120s TTL", YELLOW),
    (Inches(9.8), Inches(1.8), "Streamlit UI\nDashboard", GREEN),
]
for x, y, label, color in boxes:
    card = add_shape(slide, x, y, Inches(2.5), Inches(1.6), CARD_BG)
    add_text_box(slide, x, y + Inches(0.2), Inches(2.5), Inches(1.2),
        label, font_size=16, bold=True, color=color, align=PP_ALIGN.CENTER)

add_text_box(slide, Inches(0.8), Inches(3.8), Inches(11.3), Inches(0.3),
    "___________________________________________________________", font_size=12, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

fallback_items = [
    ("Fallback Path", ""),
    ("NTES Unavailable →", " RailSimulator generates time-aware synthetic data"),
    ("Distance Lookup →", " Station name resolver matches NTES names to known routes"),
    ("All data rendered →", " same dashboard, FALLBACK badge shown"),
]
for i, (left_text, right_text) in enumerate(fallback_items):
    y = Inches(4.2) + Inches(0.5) * i
    if i == 0:
        add_text_box(slide, Inches(1), y, Inches(11), Inches(0.4),
            left_text, font_size=20, bold=True, color=YELLOW)
    else:
        add_text_box(slide, Inches(1.5), y, Inches(3), Inches(0.4),
            left_text, font_size=16, color=TEXT_SECONDARY)
        add_text_box(slide, Inches(4.5), y, Inches(7.5), Inches(0.4),
            right_text, font_size=16, color=TEXT_MUTED)

# ── Slide 8: Analytics ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
    "Analytics Dashboard", font_size=40, bold=True, color=WHITE)
add_line(slide, Inches(0.8), Inches(1.2), Inches(3), BLUE_GRAD)
add_bullet_box(slide, Inches(1), Inches(1.8), Inches(11.3), Inches(5), [
    "Station Analysis: Average/max delay per station, train count breakdown",
    "Weather Impact: Delay rates across clear, rain, fog, and storm conditions",
    "Peak Hour Trends: 30% higher delays during morning (7-10) and evening (16-20) peaks",
    "Route Performance: Top 10 most delayed routes with delay heatmap",
    "Monthly Trends: Seasonal patterns — monsoon months (Jun-Sep) show 40% more delays",
    "Export: Download any analytics view as CSV",
], font_size=19, color=TEXT_SECONDARY)

# ── Slide 9: Route Map ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
    "Interactive Route Network", font_size=40, bold=True, color=WHITE)
add_line(slide, Inches(0.8), Inches(1.2), Inches(3), BLUE_GRAD)
add_text_box(slide, Inches(1), Inches(1.8), Inches(11.3), Inches(0.5),
    "Geospatial visualization of India's railway network", font_size=20, color=TEXT_SECONDARY)
add_shape(slide, Inches(1), Inches(2.5), Inches(11.3), Inches(4.5), CARD_BG)
add_text_box(slide, Inches(1.5), Inches(2.8), Inches(10.3), Inches(0.5),
    "30 Stations • 46 Routes • Color-coded by delay severity", font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_bullet_box(slide, Inches(1.5), Inches(3.5), Inches(10.3), Inches(3), [
    "Folium map with CartoDB dark tiles",
    "Green markers: avg delay < 8 min  |  Orange markers: 8-15 min  |  Red markers: > 15 min",
    "Line thickness and color indicate route delay severity",
    "Click any station for popup with average delay statistics",
    "Nagpur central focus with all major route connections",
], font_size=18, color=TEXT_SECONDARY)

# ── Slide 10: Key Results ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
    "Key Results & Impact", font_size=40, bold=True, color=WHITE)
add_line(slide, Inches(0.8), Inches(1.2), Inches(3), BLUE_GRAD)

results = [
    ("Real-Time Visibility", "Live train tracking with < 15s latency from NTES API"),
    ("ML-Powered Insights", "74% accurate delay prediction with 13 feature dimensions"),
    ("100k+ Records Analyzed", "Comprehensive 3-year synthetic dataset with 30 stations"),
    ("Automatic Fallback", "Time-aware simulation when NTES is unavailable"),
    ("Operator-Friendly", "Single-tab dashboard with search, filter, and CSV export"),
    ("Interactive Mapping", "Geospatial route visualization with delay heat indicators"),
]
for i, (title, desc) in enumerate(results):
    col = i % 2
    row = i // 2
    x = Inches(0.8) + Inches(6) * col
    y = Inches(1.6) + Inches(1.8) * row
    card = add_shape(slide, x, y, Inches(5.5), Inches(1.5), CARD_BG)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.1), Inches(5.1), Inches(0.4),
        f"✦ {title}", font_size=20, bold=True, color=BLUE_GRAD)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.6), Inches(5.1), Inches(0.7),
        desc, font_size=16, color=TEXT_SECONDARY)

# ── Slide 11: Future Scope ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
    "Future Scope", font_size=40, bold=True, color=WHITE)
add_line(slide, Inches(0.8), Inches(1.2), Inches(3), BLUE_GRAD)
add_bullet_box(slide, Inches(1), Inches(1.8), Inches(11.3), Inches(5), [
    "All-India Coverage: Expand from single-station (NGP) to all 8,000+ stations",
    "PNR Status Integration: Real-time passenger-specific train tracking",
    "WhatsApp / Telegram Alerts: Proactive delay notifications for passengers",
    "Arrival Prediction: ML-based ETA estimation using live GPS + historical patterns",
    "Historical Trend Dashboard: Year-over-year delay comparison and seasonal forecasting",
    "Multi-station View: Compare performance across regions, zones, and divisions",
], font_size=19, color=TEXT_SECONDARY)

# ── Slide 12: Thank You ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_shape(slide, Inches(0.8), Inches(1.8), Inches(11.7), Inches(4.2), CARD_BG)
add_line(slide, Inches(0.8), Inches(1.8), Inches(11.7), PURPLE_GRAD)
add_text_box(slide, Inches(2), Inches(2.2), Inches(9.3), Inches(1.2),
    "Thank You", font_size=60, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(slide, Inches(2), Inches(3.3), Inches(9.3), Inches(0.8),
    "RailPulse — Railway Intelligence System", font_size=28, color=BLUE_GRAD, align=PP_ALIGN.CENTER)
add_line(slide, Inches(5), Inches(4.2), Inches(3.3), PURPLE_GRAD)
add_text_box(slide, Inches(2), Inches(4.5), Inches(9.3), Inches(0.5),
    "https://github.com/aniket-diyewar/RailPulse", font_size=18, color=TEXT_SECONDARY, align=PP_ALIGN.CENTER)
add_text_box(slide, Inches(2), Inches(5.2), Inches(9.3), Inches(0.5),
    "Built with Streamlit • Scikit-learn • Plotly • NTES API", font_size=16, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

out_path = r"C:\Users\anike\Downloads\HackerEarth\RailPulse_Presentation.pptx"
prs.save(out_path)
print(f"Saved to {out_path}")
print(f"{len(prs.slides)} slides created")
